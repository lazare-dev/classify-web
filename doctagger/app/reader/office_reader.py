# Office file reader
import logging
import sys
from pathlib import Path

# Check for Office document libraries
try:
    import docx
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# For Windows file properties
try:
    import win32com.client
    import pythoncom
    HAS_WIN32COM = True
except ImportError:
    HAS_WIN32COM = False

class OfficeReader:
    """Reader for Microsoft Office files."""
    
    def read_document(self, file_path: Path) -> str:
        """Extract text from an Office document based on its file type."""
        file_ext = file_path.suffix.lower()
        
        # Office Open XML formats
        if file_ext == '.docx':
            return self._read_docx(file_path)
        elif file_ext == '.xlsx':
            return self._read_xlsx(file_path)
        elif file_ext == '.pptx':
            return self._read_pptx(file_path)
        
        # Legacy Office formats - try COM objects
        elif file_ext in ['.doc', '.xls', '.ppt'] and HAS_WIN32COM:
            return self._read_legacy_office(file_path)
        else:
            # Fall back to text reading
            try:
                with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                    content = f.read()
                    return content
            except Exception as e:
                logging.error(f"Couldn't read Office file as text: {e}")
                raise ValueError(f"Unsupported or unreadable Office file: {file_path}")
    
    def _read_docx(self, file_path: Path) -> str:
        """Read text from a DOCX file."""
        if not HAS_DOCX:
            raise ImportError("python-docx is required to read DOCX files.")
        
        try:
            doc = docx.Document(file_path)
            paragraphs = [para.text for para in doc.paragraphs]
            
            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        paragraphs.append(cell.text)
            
            text = "\n".join(paragraphs)
            return text
        except Exception as e:
            logging.error(f"Error reading DOCX: {e}")
            raise
    
    def _read_xlsx(self, file_path: Path) -> str:
        """Read text from an XLSX file."""
        if not HAS_OPENPYXL:
            raise ImportError("openpyxl is required to read XLSX files.")
        
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            
            text = []
            for sheet in workbook.worksheets:
                sheet_text = []
                for row in sheet.iter_rows():
                    row_text = []
                    for cell in row:
                        if cell.value:
                            row_text.append(str(cell.value))
                    if row_text:
                        sheet_text.append(" ".join(row_text))
                
                if sheet_text:
                    text.append(f"--- Sheet: {sheet.title} ---")
                    text.extend(sheet_text)
            
            return "\n".join(text)
        except Exception as e:
            logging.error(f"Error reading XLSX: {e}")
            raise
    
    def _read_pptx(self, file_path: Path) -> str:
        """Read text from a PPTX file."""
        if not HAS_PPTX:
            raise ImportError("python-pptx is required to read PPTX files.")
        
        try:
            presentation = Presentation(file_path)
            
            text = []
            for slide_index, slide in enumerate(presentation.slides, 1):
                slide_text = []
                
                # Extract title if available
                if slide.shapes.title and slide.shapes.title.has_text_frame:
                    title_text = slide.shapes.title.text
                    slide_text.append(f"Title: {title_text}")
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                
                if slide_text:
                    text.append(f"--- Slide {slide_index} ---")
                    text.extend(slide_text)
            
            return "\n".join(text)
        except Exception as e:
            logging.error(f"Error reading PPTX: {e}")
            raise
    
    def _read_legacy_office(self, file_path: Path) -> str:
        """Read text from legacy Office formats using COM objects."""
        if not HAS_WIN32COM:
            raise ImportError("pywin32 is required to read legacy Office files.")
        
        # Initialize COM
    def _read_legacy_office(self, file_path: Path) -> str:
        """Read text from legacy Office formats using COM objects."""
        if not HAS_WIN32COM:
            raise ImportError("pywin32 is required to read legacy Office files.")
        
        # Initialize COM
        pythoncom.CoInitialize()
        
        try:
            file_ext = file_path.suffix.lower()
            abs_path = str(file_path.resolve())
            
            if file_ext == '.doc':
                # Try Word
                try:
                    word = win32com.client.Dispatch("Word.Application")
                    word.Visible = False
                    doc = word.Documents.Open(abs_path)
                    text = doc.Content.Text
                    doc.Close()
                    word.Quit()
                    return text
                except Exception as word_err:
                    logging.error(f"Error reading DOC with Word COM: {word_err}")
                    raise
                    
            elif file_ext == '.xls':
                # Try Excel
                try:
                    excel = win32com.client.Dispatch("Excel.Application")
                    excel.Visible = False
                    excel.DisplayAlerts = False
                    workbook = excel.Workbooks.Open(abs_path)
                    
                    text = []
                    for sheet in workbook.Sheets:
                        sheet_text = []
                        used_range = sheet.UsedRange
                        for row in used_range.Rows:
                            row_text = []
                            for cell in row.Cells:
                                if cell.Value:
                                    row_text.append(str(cell.Value))
                            if row_text:
                                sheet_text.append(" ".join(row_text))
                                
                        if sheet_text:
                            text.append(f"--- Sheet: {sheet.Name} ---")
                            text.extend(sheet_text)
                    
                    workbook.Close(False)
                    excel.Quit()
                    
                    return "\n".join(text)
                except Exception as excel_err:
                    logging.error(f"Error reading XLS with Excel COM: {excel_err}")
                    raise
                    
            elif file_ext == '.ppt':
                # Try PowerPoint
                try:
                    ppt = win32com.client.Dispatch("PowerPoint.Application")
                    ppt.Visible = False
                    presentation = ppt.Presentations.Open(abs_path)
                    
                    text = []
                    for slide_idx in range(1, presentation.Slides.Count + 1):
                        slide = presentation.Slides(slide_idx)
                        text.append(f"--- Slide {slide_idx} ---")
                        
                        for shape in slide.Shapes:
                            if hasattr(shape, "TextFrame") and shape.TextFrame.HasText:
                                text.append(shape.TextFrame.TextRange.Text)
                    
                    presentation.Close()
                    ppt.Quit()
                    
                    return "\n".join(text)
                except Exception as ppt_err:
                    logging.error(f"Error reading PPT with PowerPoint COM: {ppt_err}")
                    raise
            
            raise ValueError(f"Unsupported Office file type: {file_ext}")
            
        except Exception as e:
            logging.error(f"Error reading legacy Office file: {e}")
            raise
        finally:
            pythoncom.CoUninitialize()